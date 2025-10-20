import json
import os
from typing import Dict, Any
import urllib.request
import urllib.parse
from io import BytesIO
import base64

TELEGRAM_TOKEN = "8310137961:AAFYWzx442w3bx3Y7VPvuMDDbyl0W576qbQ"
TELEGRAM_API = f"https://api.telegram.org/bot{TELEGRAM_TOKEN}"

def handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    '''
    Business: Telegram бот для создания презентаций PPTX
    Args: event - webhook данные от Telegram
          context - контекст выполнения функции
    Returns: HTTP response для Telegram
    '''
    method: str = event.get('httpMethod', 'POST')
    
    if method == 'OPTIONS':
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'GET, POST, OPTIONS',
                'Access-Control-Allow-Headers': 'Content-Type',
                'Access-Control-Max-Age': '86400'
            },
            'body': '',
            'isBase64Encoded': False
        }
    
    if method == 'POST':
        try:
            body_data = json.loads(event.get('body', '{}'))
            
            if 'message' in body_data:
                message = body_data['message']
                chat_id = message['chat']['id']
                text = message.get('text', '')
                
                if text.startswith('/start'):
                    send_start_message(chat_id)
                elif text.startswith('/help'):
                    send_help_message(chat_id)
                elif text.startswith('/examples'):
                    send_examples_message(chat_id)
                else:
                    create_and_send_presentation(chat_id, text)
        except Exception as e:
            pass
        
        return {
            'statusCode': 200,
            'headers': {'Content-Type': 'application/json'},
            'body': json.dumps({'ok': True}),
            'isBase64Encoded': False
        }
    
    return {
        'statusCode': 200,
        'headers': {'Content-Type': 'text/plain'},
        'body': 'Presentation Bot is running!',
        'isBase64Encoded': False
    }

def send_message(chat_id: int, text: str, parse_mode: str = 'HTML') -> None:
    params = {
        'chat_id': str(chat_id),
        'text': text,
        'parse_mode': parse_mode
    }
    
    try:
        url = f"{TELEGRAM_API}/sendMessage?" + urllib.parse.urlencode(params)
        req = urllib.request.Request(url)
        
        with urllib.request.urlopen(req) as response:
            response.read()
    except Exception as e:
        pass

def send_start_message(chat_id: int) -> None:
    text = """📊 <b>Добро пожаловать в Presentation AI Bot!</b>

Я создаю настоящие презентации в формате PPTX и отправляю тебе готовый файл!

<b>🎯 Что я умею:</b>
• Генерировать презентации по любой теме
• Создавать профессиональные слайды
• Добавлять структуру и оформление
• Отправлять готовый PPTX файл

<b>🚀 Как использовать:</b>
Просто опиши тему презентации и я создам её!

<b>Примеры:</b>
"Презентация про здоровое питание на 5 слайдов"
"Питч для стартапа по доставке еды, 7 слайдов"
"Обучение сотрудников по безопасности, 10 слайдов"

<b>📋 Команды:</b>
/help - Подробная инструкция
/examples - Примеры презентаций

Напиши тему своей презентации и получи готовый файл! ✨"""
    
    send_message(chat_id, text)

def send_help_message(chat_id: int) -> None:
    text = """❓ <b>Как создать презентацию</b>

<b>📝 Опиши свою презентацию:</b>

<b>1️⃣ Тема (обязательно)</b>
О чём презентация?
• "Презентация про маркетинг"
• "Защита диплома по экономике"
• "Питч для инвесторов"

<b>2️⃣ Количество слайдов</b>
От 3 до 15 слайдов
• "5 слайдов"
• "10 слайдов"

<b>3️⃣ Стиль (необязательно)</b>
• "Деловой стиль"
• "Креативный дизайн"
• "Минималистичный"

<b>✅ Примеры запросов:</b>

💼 "Презентация о нашей компании, 7 слайдов, деловой стиль"

📚 "Урок по истории России, 10 слайдов, для школьников"

🚀 "Питч стартапа по AI, 5 слайдов, современный дизайн"

<b>📊 Что будет в презентации:</b>
• Титульный слайд с названием
• Введение с описанием темы
• Основные слайды по теме
• Заключение с выводами

<b>⏱ Время создания:</b>
10-15 секунд

Опиши свою тему и получи готовый PPTX файл! 🎬"""
    
    send_message(chat_id, text)

def send_examples_message(chat_id: int) -> None:
    text = """📊 <b>Примеры презентаций</b>

<b>Для бизнеса:</b>
💼 "Презентация компании для клиентов, 8 слайдов"
💼 "Квартальный отчёт для руководства, 10 слайдов"
💼 "Питч для инвесторов, 5 слайдов"

<b>Для образования:</b>
📚 "Урок по физике: законы Ньютона, 7 слайдов"
📚 "Дипломная защита по программированию, 12 слайдов"
📚 "Проект для школы про экологию, 6 слайдов"

<b>Для маркетинга:</b>
📈 "Презентация нового продукта, 8 слайдов"
📈 "Маркетинговая стратегия на год, 10 слайдов"
📈 "Анализ конкурентов, 6 слайдов"

<b>Для стартапов:</b>
🚀 "Питч-дек для акселератора, 10 слайдов"
🚀 "Бизнес-план мобильного приложения, 12 слайдов"
🚀 "Презентация MVP инвесторам, 7 слайдов"

<b>Креативные:</b>
🎨 "Портфолио дизайнера, 8 слайдов"
🎨 "Идея для нового проекта, 5 слайдов"
🎨 "План мероприятия, 6 слайдов"

Выбери пример или напиши свою тему! 💡"""
    
    send_message(chat_id, text)

def create_and_send_presentation(chat_id: int, user_text: str) -> None:
    send_message(chat_id, "⏳ Создаю презентацию... Это займёт 10-15 секунд")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        topic = extract_topic(user_text)
        num_slides = extract_slide_count(user_text)
        
        add_title_slide(prs, topic)
        add_content_slides(prs, topic, num_slides - 2)
        add_conclusion_slide(prs, topic)
        
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        send_document(chat_id, pptx_io, f"{topic[:30]}.pptx")
        send_message(chat_id, f"✅ <b>Готово!</b>\n\nПрезентация '{topic}' создана!\nСлайдов: {num_slides}\n\nМожешь скачать и редактировать в PowerPoint, Google Slides или Keynote 📊")
        
    except ImportError:
        send_message(chat_id, """📊 <b>Создание презентации</b>

Твоя тема: "{}"

<b>Структура презентации:</b>

<b>Слайд 1: Титульный</b>
📌 {}

<b>Слайд 2: Введение</b>
• Актуальность темы
• Цели и задачи
• Краткий обзор

<b>Слайды 3-{}: Основная часть</b>
• Ключевые пункты по теме
• Факты и статистика
• Примеры и иллюстрации
• Детальный анализ

<b>Слайд {}: Заключение</b>
• Выводы
• Рекомендации
• Контакты

<b>💡 Рекомендации по оформлению:</b>
✓ Используй единый стиль
✓ Крупный шрифт (24+ pt)
✓ Минимум текста на слайде
✓ Качественные изображения
✓ Инфографика для данных

<b>🛠 Создай в:</b>
• Google Slides (бесплатно онлайн)
• PowerPoint
• Canva (готовые шаблоны)

Удачи с презентацией! 🚀""".format(
            user_text,
            extract_topic(user_text),
            extract_slide_count(user_text) - 1,
            extract_slide_count(user_text)
        ))

def extract_topic(text: str) -> str:
    text_clean = text.lower()
    for word in ['презентация про', 'презентация о', 'презентация', 'про ', 'о ', 'на тему']:
        text_clean = text_clean.replace(word, '')
    
    for word in [' слайд', ' на ', 'деловой', 'креативный', 'минималист']:
        if word in text_clean:
            text_clean = text_clean.split(word)[0]
    
    return text_clean.strip().capitalize()[:50] or "Презентация"

def extract_slide_count(text: str) -> int:
    import re
    numbers = re.findall(r'\d+', text)
    if numbers:
        num = int(numbers[0])
        return max(3, min(num, 15))
    return 5

def add_title_slide(prs, topic):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(68, 84, 106)
    p.alignment = PP_ALIGN.CENTER

def add_content_slides(prs, topic, count):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    content_points = generate_content_points(topic, count)
    
    for i, point in enumerate(content_points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{point['title']}"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
        content_frame = content_box.text_frame
        
        for bullet in point['bullets']:
            p = content_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(89, 89, 89)

def add_conclusion_slide(prs, topic):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Заключение"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(68, 84, 106)
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4.5))
    content_frame = content_box.text_frame
    
    conclusions = [
        f"Мы рассмотрели ключевые аспекты темы '{topic}'",
        "Проанализировали основные концепции и подходы",
        "Изучили практическое применение",
        "Спасибо за внимание!"
    ]
    
    for conclusion in conclusions:
        p = content_frame.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(89, 89, 89)

def generate_content_points(topic, count):
    points = []
    base_sections = [
        {"title": "Введение", "bullets": [
            f"Актуальность темы '{topic}'",
            "Цели и задачи презентации",
            "Краткий обзор содержания"
        ]},
        {"title": "Основные понятия", "bullets": [
            "Определения и термины",
            "Теоретические основы",
            "Ключевые концепции"
        ]},
        {"title": "Анализ", "bullets": [
            "Текущее состояние вопроса",
            "Статистика и данные",
            "Тренды и тенденции"
        ]},
        {"title": "Практическое применение", "bullets": [
            "Примеры использования",
            "Кейсы и исследования",
            "Лучшие практики"
        ]},
        {"title": "Преимущества", "bullets": [
            "Ключевые выгоды",
            "Положительные эффекты",
            "Долгосрочная перспектива"
        ]},
        {"title": "Вызовы и решения", "bullets": [
            "Текущие проблемы",
            "Возможные риски",
            "Способы преодоления"
        ]},
        {"title": "Выводы", "bullets": [
            "Основные результаты",
            "Практические рекомендации",
            "Дальнейшие шаги"
        ]}
    ]
    
    for i in range(min(count, len(base_sections))):
        points.append(base_sections[i])
    
    return points

def send_document(chat_id: int, file_data: BytesIO, filename: str) -> None:
    boundary = '----WebKitFormBoundary7MA4YWxkTrZu0gW'
    
    body = (
        f'--{boundary}\r\n'
        f'Content-Disposition: form-data; name="chat_id"\r\n\r\n'
        f'{chat_id}\r\n'
        f'--{boundary}\r\n'
        f'Content-Disposition: form-data; name="document"; filename="{filename}"\r\n'
        f'Content-Type: application/vnd.openxmlformats-officedocument.presentationml.presentation\r\n\r\n'
    ).encode('utf-8')
    
    body += file_data.read()
    body += f'\r\n--{boundary}--\r\n'.encode('utf-8')
    
    try:
        req = urllib.request.Request(
            f"{TELEGRAM_API}/sendDocument",
            data=body,
            headers={'Content-Type': f'multipart/form-data; boundary={boundary}'}
        )
        
        with urllib.request.urlopen(req) as response:
            response.read()
    except Exception as e:
        pass
